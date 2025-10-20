import os
import json
import asyncio
import time
import openai
import httpx
from dotenv import load_dotenv
from semantic_kernel import Kernel
from semantic_kernel.connectors.ai.open_ai import AzureChatCompletion
from semantic_kernel.functions import kernel_function
from pptx import Presentation
from pptx.util import Inches, Pt
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt
from semantic_kernel.connectors.ai.open_ai import OpenAITextPromptExecutionSettings
from semantic_kernel.contents.chat_history import ChatHistory

# Load environment variables from .env file
load_dotenv()

class ContentGenerationPlugin:
    @kernel_function(
        name="generate_document_content",
        description="Generate content for a PPTX or DOCX document based on a topic and document type."
    )
    async def generate_content(self, topic: str, doc_type: str, num_sections: int = 5) -> dict:
        """
        Generate content for a document using Azure OpenAI, tailored for PPTX or DOCX.
        
        :param topic: The topic of the document.
        :param doc_type: The document type ('pptx' or 'docx').
        :param num_sections: Number of sections (slides for PPTX, paragraphs for DOCX).
        :return: Dictionary with section titles and content.
        """
        kernel = Kernel()
        kernel.add_service(
            AzureChatCompletion(
                deployment_name=os.getenv("AZURE_OPENAI_MODEL"),
                endpoint=os.getenv("AZURE_OPENAI_ENDPOINT"),
                api_key=os.getenv("AZURE_OPENAI_API_KEY"),
                api_version="2023-05-15"
            )
        )
        
        # Fine-tune prompt with explicit JSON formatting instructions
        if doc_type.lower() == "pptx":
            prompt = f"""
            You are an expert content generator. Generate content for a PowerPoint presentation on the topic: '{topic}'.
            Create exactly {num_sections} slides.
            For each slide, provide:
            - A short, catchy title (max 5-7 words, suitable for a slide).
            - A list of 3-4 concise bullet points (max 10 words each, presentation-friendly).
            
            Return the output in valid JSON format, wrapped in triple backticks:
            ```json
            {{
                "sections": [
                    {{"title": "Slide Title", "content": ["bullet1", "bullet2", ...]}},
                    ...
                ]
            }}
            ```
            Ensure the JSON is properly formatted and contains no extra text outside the JSON structure.
            """
        else:  # docx
            prompt = f"""
            You are an expert content generator. Generate content for a Word document on the topic: '{topic}'.
            Create exactly {num_sections} sections.
            For each section, provide:
            - A clear, descriptive heading (max 10-12 words, suitable for a document).
            - A list of 4-5 detailed bullet points (max 20 words each, suitable for a report).
            
            Return the output in valid JSON format, wrapped in triple backticks:
            ```json
            {{
                "sections": [
                    {{"title": "Section Heading", "content": ["bullet1", "bullet2", ...]}},
                    ...
                ]
            }}
            ```
            Ensure the JSON is properly formatted and contains no extra text outside the JSON structure.
            """
        
        settings = OpenAITextPromptExecutionSettings(
            max_tokens=1000,
            temperature=0.7
        )
        
        chat_service = kernel.get_service(type=AzureChatCompletion)
        
        # Build a ChatHistory and add the user prompt
        chat_history = ChatHistory()
        chat_history.add_user_message(prompt)

        # Call the chat service with retries/backoff to handle transient timeouts
        max_retries = 3
        backoff_seconds = 1
        last_exc = None
        for attempt in range(1, max_retries + 1):
            try:
                response = await chat_service.get_chat_message_content(
                    chat_history=chat_history,
                    settings=settings,
                )
                last_exc = None
                break
            except Exception as exc:
                last_exc = exc
                # Detect likely timeout/connection exceptions to retry
                is_timeout = isinstance(exc, openai.APITimeoutError) or isinstance(exc, httpx.TimeoutException) or 'timed out' in str(exc).lower()
                print(f"Attempt {attempt} failed: {exc}")
                if attempt < max_retries and is_timeout:
                    wait = backoff_seconds * (2 ** (attempt - 1))
                    print(f"Retrying in {wait} seconds...")
                    await asyncio.sleep(wait)
                    continue
                # Not a retryable error or out of attempts: re-raise
                raise
        if last_exc is not None:
            # If we fell through with an exception, raise a clearer error
            raise last_exc
        
        if not response or not response.content:
            print("Error: Empty response from LLM.")
            raise ValueError("LLM returned an empty response.")
        
        # Log raw response for debugging
        #print(f"Raw LLM response: {response.content}")
        
        # Extract JSON from response (strip triple backticks if present)
        content = response.content.strip()
        if content.startswith("```json") and content.endswith("```"):
            content = content[7:-3].strip()
        
        try:
            return json.loads(content)
        except json.JSONDecodeError as e:
            print(f"Error: Failed to parse JSON from LLM response: {content}")
            raise ValueError(f"Failed to parse JSON from LLM response: {str(e)}")

class PptxGenerationPlugin:
    @kernel_function(
        name="create_pptx",
        description="Create a PowerPoint presentation from generated content."
    )
    def create_pptx(self, sections_content: list, output_file: str = "generated_ppt.pptx") -> None:
        """
        Create a PowerPoint presentation using python-pptx.
        
        :param sections_content: List of dictionaries with 'title' and 'content'.
        :param output_file: Path to save the PPTX.
        """
        prs = Presentation()
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = "Generated Presentation"
        subtitle = slide.placeholders[1]
        subtitle.text = "Powered by Azure OpenAI & Semantic Kernel"
        
        # Add content slides
        bullet_slide_layout = prs.slide_layouts[1]
        for section in sections_content:
            slide = prs.slides.add_slide(bullet_slide_layout)
            title = slide.shapes.title
            title.text = section["title"]
            
            body_shape = slide.shapes.placeholders[1]
            tf = body_shape.text_frame
            for bullet in section["content"]:
                p = tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(18)
        
        # Ensure the directory exists (only if a directory component was provided)
        #dirpath = os.path.dirname(output_file)
        dirpath = os.getenv("LOCAL_PPTX_DIRECTORY")
        if dirpath:
            os.makedirs(dirpath, exist_ok=True)
        # Construct the full path to save the document
        full_path = os.path.join(dirpath, output_file)
        prs.save(full_path)
        print(f"PPTX saved locally to {os.path.abspath(full_path)}")

class DocxGenerationPlugin:
    @kernel_function(
        name="create_docx",
        description="Create a Word document from generated content."
    )
    def create_docx(self, sections_content: list, output_file: str = "generated_doc.docx") -> None:
        """
        Create a Word document using python-docx.
        
        :param sections_content: List of dictionaries with 'title' and 'content'.
        :param output_file: Path to save the DOCX.
        """
        doc = Document()
        
        # Add title
        doc.add_heading("Generated Document", 0)
        doc.add_paragraph("Powered by Azure OpenAI & Semantic Kernel", style="Subtitle")
        
        # Add content sections
        for section in sections_content:
            doc.add_heading(section["title"], level=1)
            for bullet in section["content"]:
                doc.add_paragraph(bullet, style="List Bullet")
        
        # Ensure the directory exists (only if a directory component was provided)
        #dirpath = os.path.dirname(output_file)
        dirpath = os.getenv("LOCAL_DOCX_DIRECTORY")
        if dirpath:
            os.makedirs(dirpath, exist_ok=True)
        # Construct the full path to save the document
        full_path = os.path.join(dirpath, output_file)
        doc.save(full_path)
        print(f"DOCX saved locally to {os.path.abspath(full_path)}")

async def main():
    # Initialize Semantic Kernel
    kernel = Kernel()
    kernel.add_service(
        AzureChatCompletion(
            deployment_name=os.getenv("AZURE_OPENAI_MODEL"),
            endpoint=os.getenv("AZURE_OPENAI_ENDPOINT"),
            api_key=os.getenv("AZURE_OPENAI_API_KEY"),
            api_version="2023-05-15"
        )
    )
    kernel.add_plugin(ContentGenerationPlugin(), plugin_name="ContentGeneration")
    kernel.add_plugin(PptxGenerationPlugin(), plugin_name="PptxGeneration")
    kernel.add_plugin(DocxGenerationPlugin(), plugin_name="DocxGeneration")

    # Step 1: Ask for topic
    print("What topic would you like for your document?")
    topic = input().strip()
    if not topic:
        print("Topic cannot be empty. Using default: 'The Future of AI'")
        topic = "The Future of AI"

    # Step 2: Ask for document type
    while True:
        print("Which document type do you want? (Enter 'pptx' or 'docx')")
        doc_type = input().strip().lower()
        if doc_type in ["pptx", "docx"]:
            break
        print("Invalid input. Please enter 'pptx' or 'docx'.")

    # Step 3: Ask for filename
    print("Enter the filename (e.g., 'my_document.pptx' or 'my_document.docx'). Leave blank for default.")
    filename = input().strip()
    if not filename:
        filename = f"generated_{'ppt' if doc_type == 'pptx' else 'doc'}.{doc_type}"
    else:
        # Ensure correct extension
        if not filename.endswith(f".{doc_type}"):
            filename = f"{filename}.{doc_type}"

    # Generate content using Semantic Kernel
    content_plugin = kernel.get_plugin("ContentGeneration")
    result = await kernel.invoke(content_plugin["generate_document_content"], topic=topic, doc_type=doc_type, num_sections=5)
    sections_content = result.value["sections"]

    # Create and save the document
    if doc_type == "pptx":
        pptx_plugin = kernel.get_plugin("PptxGeneration")
        await kernel.invoke(pptx_plugin["create_pptx"], sections_content=sections_content, output_file=filename)
    else:
        docx_plugin = kernel.get_plugin("DocxGeneration")
        await kernel.invoke(docx_plugin["create_docx"], sections_content=sections_content, output_file=filename)

if __name__ == "__main__":
    asyncio.run(main())