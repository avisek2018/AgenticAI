"""
Multiple Function Tools (Interactive Demo)

This demo shows an agent with MULTIPLE tools:

The agent automatically chooses the right tool based on your question.
"""

import asyncio
import os
from typing import Annotated
from pydantic import Field
from dotenv import load_dotenv
from datetime import datetime
import requests

from agent_framework.azure import AzureOpenAIChatClient

from agent_framework import ChatAgent
from agent_framework.azure import AzureAIAgentClient
from azure.identity.aio import AzureCliCredential

from agent_framework.devui import serve


import json
from typing import Annotated
from pydantic import Field
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt


from azure.ai.translation.text import TextTranslationClient
from azure.core.credentials import AzureKeyCredential
from azure.ai.textanalytics import TextAnalyticsClient

# Load environment variables
load_dotenv()

# ENDPOINT = os.getenv("AZURE_OPENAI_ENDPOINT")
# DEPLOYMENT = os.getenv("AZURE_OPENAI_CHAT_DEPLOYMENT_NAME")
# API_KEY = os.getenv("AZURE_OPENAI_API_KEY")
# API_VERSION = os.getenv("AZURE_OPENAI_API_VERSION", "2024-07-01-preview")


PROJECT_ENDPOINT = os.getenv("AZURE_AI_PROJECT_ENDPOINT")
AGENT_ID = os.getenv("AZURE_AI_AGENT_ID")

AZURE_TRANSLATOR_ENDPOINT = os.getenv("AZURE_TRANSLATOR_ENDPOINT", "https://api.cognitive.microsofttranslator.com/")
AZURE_TRANSLATOR_API_KEY = os.getenv("AZURE_TRANSLATOR_API_KEY")
AZURE_TRANSLATOR_REGION = os.getenv("AZURE_TRANSLATOR_REGION", "eastus")

# Global Translator Client (initialized once for efficiency)
translator_client = None
if AZURE_TRANSLATOR_API_KEY and AZURE_TRANSLATOR_REGION:
    translator_client = TextTranslationClient(
        endpoint=AZURE_TRANSLATOR_ENDPOINT,
        credential=AzureKeyCredential(AZURE_TRANSLATOR_API_KEY),
        region=AZURE_TRANSLATOR_REGION
    )

# Azure Text Analytics setup
TEXT_ANALYTICS_ENDPOINT = os.getenv("AZURE_LANGUAGE_SERVICE_ENDPOINT")
TEXT_ANALYTICS_KEY = os.getenv("AZURE_LANGUAGE_SERVICE_KEY")
text_analytics_client = None
if TEXT_ANALYTICS_ENDPOINT and TEXT_ANALYTICS_KEY:
    text_analytics_client = TextAnalyticsClient(
        endpoint=TEXT_ANALYTICS_ENDPOINT,
        credential=AzureKeyCredential(TEXT_ANALYTICS_KEY)
    )


# Tool 1: PPTX Generation
def generate_pptx(
    content_json: Annotated[str, Field(description="JSON string containing sections with titles and content for the PPTX")],
    filename: Annotated[str, Field(description="Filename to save the PPTX (e.g., 'output.pptx')")] = "generated_ppt.pptx"
) -> str:
    """Generate a PowerPoint presentation from provided JSON content."""
    try:
        # Parse JSON content
        try:
            sections_content = json.loads(content_json)["sections"]
        except json.JSONDecodeError as e:
            return f"Error: Failed to parse JSON content for PPTX: {str(e)}"

        # Validate sections content
        if not isinstance(sections_content, list) or not all(
            isinstance(section, dict) and "title" in section and "content" in section
            for section in sections_content
        ):
            return "Error: Invalid JSON structure. Expected 'sections' with 'title' and 'content' fields."

        # Create PPTX
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = "Generated Presentation"
        subtitle = slide.placeholders[1]
        subtitle.text = "Powered by Azure AI"

        bullet_slide_layout = prs.slide_layouts[1]
        for section in sections_content:
            slide = prs.slides.add_slide(bullet_slide_layout)
            title = slide.shapes.title
            title.text = section["title"]
            body_shape = slide.shapes.placeholders[1]
            tf = body_shape.text_frame
            for bullet in section["content"]:
                p = tf.add_paragraph()
                p.text = str(bullet)
                p.font.size = Pt(18)

        # Ensure correct extension and directory
        if not filename.endswith(".pptx"):
            filename = f"{filename}.pptx"
        os.makedirs(os.path.dirname(filename) if os.path.dirname(filename) else ".", exist_ok=True)
        prs.save(filename)
        return f"PPTX saved to {os.path.abspath(filename)}"
    except Exception as e:
        return f"Error generating PPTX: {str(e)}"

# Tool 2: DOCX Generation
def generate_docx(
    content_json: Annotated[str, Field(description="JSON string containing sections with titles and content for the DOCX")],
    filename: Annotated[str, Field(description="Filename to save the DOCX (e.g., 'output.docx')")] = "generated_doc.docx"
) -> str:
    """Generate a Word document from provided JSON content."""
    try:
        # Parse JSON content
        try:
            sections_content = json.loads(content_json)["sections"]
        except json.JSONDecodeError as e:
            return f"Error: Failed to parse JSON content for DOCX: {str(e)}"

        # Validate sections content
        if not isinstance(sections_content, list) or not all(
            isinstance(section, dict) and "title" in section and "content" in section
            for section in sections_content
        ):
            return "Error: Invalid JSON structure. Expected 'sections' with 'title' and 'content' fields."

        # Create DOCX
        doc = Document()
        doc.add_heading("Generated Document", 0)
        doc.add_paragraph("Powered by Azure AI", style="Subtitle")

        for section in sections_content:
            doc.add_heading(section["title"], level=1)
            for bullet in section["content"]:
                doc.add_paragraph(str(bullet), style="List Bullet")

        # Ensure correct extension and directory
        if not filename.endswith(".docx"):
            filename = f"{filename}.docx"
        os.makedirs(os.path.dirname(filename) if os.path.dirname(filename) else ".", exist_ok=True)
        doc.save(filename)
        return f"DOCX saved to {os.path.abspath(filename)}"
    except Exception as e:
        return f"Error generating DOCX: {str(e)}"


# Tool 3: Language Translation (inspired by TranslationPlugin in language_agent_mcp.py)
def translate_to_english(
    text: Annotated[str, Field(description="The text to translate to English")]
) -> str:
    """Translate the input text to English using Azure Translator."""
    if not translator_client:
        return "Error: Azure Translator client not initialized. Check environment variables."
    try:
        result = translator_client.translate(
            body=[{"text": text}],
            to_language=["en"]
        )
        if result and result[0].translations:
            translation = result[0].translations[0].text
            return f"Translated to English: {translation}"
        return "Unable to translate text."
    except Exception as e:
        return f"Error translating text: {str(e)}"


# Tool 4: Entity Extraction (Azure Text Analytics)
def extract_entities(
    document: Annotated[str, Field(description="Path to a text file to extract entities from, or the text content itself.")]
) -> str:
    """Extract named entities from a document or text using Azure Text Analytics."""
    if not text_analytics_client:
        return "Error: Azure Text Analytics client not initialized. Check environment variables."
    # Try to read file if path exists, else treat as text
    if os.path.exists(document):
        try:
            with open(document, "r", encoding="utf-8") as f:
                text = f.read()
        except Exception as e:
            return f"Error reading file: {e}"
    else:
        text = document
    try:
        response = text_analytics_client.recognize_entities(documents=[text])
        result_lines = []
        for doc in response:
            if not doc.is_error:
                result_lines.append("Entities found:")
                for entity in doc.entities:
                    result_lines.append(f"  Entity: {entity.text}, Category: {entity.category}, SubCategory: {getattr(entity, 'subcategory', None) or 'N/A'}")
            else:
                result_lines.append(f"Error processing document: {doc.error.message}")
        return "\n".join(result_lines) if result_lines else "No entities found."
    except Exception as e:
        return f"Error extracting entities: {e}"


async def main():
    """Interactive demo: Agent with multiple tools (async).

    This function creates the agent and launches the DevUI in a background
    thread to avoid nested event loops.
    """

    # Load system instructions from file if available
    def _load_instructions() -> str:
        candidates = [
            os.path.join(os.getcwd(), "system_prompts.txt"),
            os.path.join(os.path.dirname(__file__), "system_prompts.txt"),
        ]
        for path in candidates:
            try:
                if os.path.exists(path):
                    with open(path, "r", encoding="utf-8") as fh:
                        content = fh.read().strip()
                        if content:
                            return content
            except Exception:
                # Ignore read errors and continue to next candidate
                pass
        # Fallback default instructions
        return "You are a helpful assistant with access to multiple tools: weather, calculator, and time zone."

    INSTRUCTIONS = _load_instructions()

    # Use the async AzureCliCredential as an async context manager
    async with AzureCliCredential() as credential, ChatAgent(
        chat_client=AzureAIAgentClient(
            async_credential=credential,
            project_endpoint=PROJECT_ENDPOINT,
            agent_id=AGENT_ID,
        ),
        instructions=INSTRUCTIONS,
        tools=[generate_pptx, generate_docx, translate_to_english, extract_entities],
    ) as agent:
        

        # Start DevUI in a daemon thread because serve() runs uvicorn and may
        # call blocking event loop functions. Running in a thread prevents
        # conflicts with the running asyncio event loop here.
        import threading

        def _start_devui():
            try:
                serve(entities=[agent], port=8000, auto_open=True)
            except Exception as e:
                print(f"DevUI serve() failed: {e}")

        threading.Thread(target=_start_devui, daemon=True).start()

        # Keep the program running so the daemon thread stays alive and the
        # async agent can be interacted with. Replace this with real
        # interaction in a complete app.
        print("Agent and DevUI started. Press Ctrl+C to exit.")
        try:
            # Sleep indefinitely while event loop runs
            await asyncio.Event().wait()
        except asyncio.CancelledError:
            pass
   
        
if __name__ == "__main__":
    # Run the async main properly
    try:
        asyncio.run(main())
    except KeyboardInterrupt:
        print("Shutting down")
